from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("outputs/industrial-knowledge-brain/AssetIQ-Brain-Pitch-Deck.pptx")

slides = [
    (
        "AssetIQ Brain",
        "Unified Asset & Operations Brain for industrial knowledge intelligence.",
        ["Hackathon prototype", "Industrial Intelligence / Document Management / Knowledge Engineering / Quality"],
    ),
    (
        "The Problem",
        "Plant knowledge is scattered across disconnected systems.",
        ["Engineers lose time searching for existing information.", "Maintenance teams make decisions without complete equipment history.", "Retiring experts take undocumented operational knowledge with them."],
    ),
    (
        "The Solution",
        "An AI-powered platform that makes drawings, work orders, SOPs, inspections, manuals, and compliance evidence queryable together.",
        ["Universal ingestion", "Industrial entity extraction", "Knowledge graph", "Cited copilot", "RCA and compliance agents"],
    ),
    (
        "Prototype Demo",
        "The demo uses a cooling water pump scenario centered on P-204.",
        ["Ingests six representative document types.", "Links P-204 to work orders, P&ID, OEM manual, SOP, inspection, and OISD checklist.", "Answers operational questions with confidence and citations."],
    ),
    (
        "Technical Architecture",
        "Document AI plus industrial ontology plus hybrid RAG.",
        ["OCR and layout parsing", "NER for equipment tags, dates, failure modes, parameters, and regulations", "Graph traversal for cross-document reasoning", "Agent workflows for field actions"],
    ),
    (
        "Business Impact",
        "Faster decisions, fewer repeat failures, stronger audit readiness.",
        ["Reduce time-to-answer across functions.", "Expose hidden failure patterns before downtime.", "Generate evidence packs from linked source documents.", "Retain expert knowledge as reusable graph context."],
    ),
    (
        "Judging Criteria Fit",
        "Built for measurable evaluation.",
        ["Innovation: graph-backed industrial agents", "Business impact: downtime, compliance, knowledge retention", "Technical excellence: extraction, retrieval, and workflow reasoning", "UX: field-ready copilot and command center"],
    ),
    (
        "Next Steps",
        "Move from demo intelligence to production intelligence.",
        ["Add real PDF, image, and spreadsheet upload.", "Connect graph database and vector store.", "Use plant-specific labelled data for extraction evaluation.", "Integrate CMMS, QMS, historian, and document vault connectors."],
    ),
]


def add_textbox(slide, left, top, width, height, text, size=22, bold=False, color=RGBColor(23, 32, 38)):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

for index, (title, subtitle, bullets) in enumerate(slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(244, 247, 248)

    band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(15, 124, 128)
    band.line.fill.background()

    add_textbox(slide, Inches(0.7), Inches(0.55), Inches(11.8), Inches(0.45), "AssetIQ Brain", 13, True, RGBColor(15, 124, 128))
    add_textbox(slide, Inches(0.7), Inches(1.15), Inches(7.6), Inches(0.85), title, 34, True)
    add_textbox(slide, Inches(0.72), Inches(2.02), Inches(8.8), Inches(0.75), subtitle, 19, False, RGBColor(80, 98, 107))

    card = slide.shapes.add_shape(1, Inches(0.7), Inches(3.0), Inches(7.7), Inches(3.55))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(215, 224, 228)

    tf = card.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.25)
    tf.clear()
    for bullet_index, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if bullet_index == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(23, 32, 38)
        p.space_after = Pt(12)
        p.level = 0

    visual = slide.shapes.add_shape(1, Inches(9.0), Inches(1.35), Inches(3.55), Inches(5.2))
    visual.fill.solid()
    visual.fill.fore_color.rgb = RGBColor(24, 48, 56)
    visual.line.fill.background()

    for n, label in enumerate(["Docs", "OCR", "Graph", "Agents"]):
        y = Inches(1.75 + n * 0.95)
        node = slide.shapes.add_shape(9, Inches(9.55), y, Inches(2.35), Inches(0.52))
        node.fill.solid()
        node.fill.fore_color.rgb = RGBColor(224, 243, 241) if n == index % 4 else RGBColor(255, 255, 255)
        node.line.fill.background()
        tf_node = node.text_frame
        tf_node.text = label
        tf_node.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf_node.paragraphs[0].runs[0].font.size = Pt(16)
        tf_node.paragraphs[0].runs[0].font.bold = True
        tf_node.paragraphs[0].runs[0].font.color.rgb = RGBColor(23, 32, 38)

    add_textbox(slide, Inches(9.45), Inches(5.8), Inches(2.6), Inches(0.35), f"0{index + 1}", 16, True, RGBColor(155, 217, 213))

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(OUT)
